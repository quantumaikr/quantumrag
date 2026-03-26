"""Office document parsers: DOCX, PPTX, XLSX."""

from __future__ import annotations

from pathlib import Path

from quantumrag.core.errors import ParseError
from quantumrag.core.logging import get_logger
from quantumrag.core.models import Document, DocumentMetadata, SourceType, Table

logger = get_logger(__name__)


class DocxParser:
    """Parser for Microsoft Word (.docx) files.

    Uses python-docx (lazy import, optional dependency).
    """

    @property
    def supported_extensions(self) -> list[str]:
        return [".docx"]

    def parse(self, file_path: Path) -> Document:
        """Parse a DOCX file."""
        if not file_path.exists():
            raise ParseError(f"File not found: {file_path}", file_path=str(file_path))

        try:
            import docx
        except ImportError:
            raise ParseError(
                "python-docx is required to parse DOCX files",
                file_path=str(file_path),
                suggestion="Install it with: pip install python-docx",
            ) from None

        try:
            doc = docx.Document(str(file_path))
        except Exception as e:
            raise ParseError(
                f"Failed to open DOCX file: {e}",
                file_path=str(file_path),
            ) from e

        # Extract text from paragraphs
        paragraphs: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        content = "\n\n".join(paragraphs)

        # Extract core properties
        props = doc.core_properties
        title = props.title or file_path.stem
        author = props.author or None

        return Document(
            content=content,
            metadata=DocumentMetadata(
                source_type=SourceType.FILE,
                title=title,
                author=author,
                custom={"format": "docx"},
            ),
            raw_bytes=file_path.read_bytes(),
        )


class PptxParser:
    """Parser for Microsoft PowerPoint (.pptx) files.

    Uses python-pptx (lazy import, optional dependency).
    """

    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx"]

    def parse(self, file_path: Path) -> Document:
        """Parse a PPTX file, extracting text from slides."""
        if not file_path.exists():
            raise ParseError(f"File not found: {file_path}", file_path=str(file_path))

        try:
            from pptx import Presentation
        except ImportError:
            raise ParseError(
                "python-pptx is required to parse PPTX files",
                file_path=str(file_path),
                suggestion="Install it with: pip install python-pptx",
            ) from None

        try:
            prs = Presentation(str(file_path))
        except Exception as e:
            raise ParseError(
                f"Failed to open PPTX file: {e}",
                file_path=str(file_path),
            ) from e

        slides_text: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts: list[str] = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
            if len(parts) > 1:  # More than just the header
                slides_text.append("\n".join(parts))

        content = "\n\n".join(slides_text)

        return Document(
            content=content,
            metadata=DocumentMetadata(
                source_type=SourceType.FILE,
                title=file_path.stem,
                custom={
                    "format": "pptx",
                    "slide_count": str(len(prs.slides)),
                },
            ),
            raw_bytes=file_path.read_bytes(),
        )


class XlsxParser:
    """Parser for Microsoft Excel (.xlsx) files.

    Uses openpyxl (lazy import, optional dependency).
    """

    @property
    def supported_extensions(self) -> list[str]:
        return [".xlsx"]

    def parse(self, file_path: Path) -> Document:
        """Parse an XLSX file, converting sheets to structured text."""
        if not file_path.exists():
            raise ParseError(f"File not found: {file_path}", file_path=str(file_path))

        try:
            import openpyxl
        except ImportError:
            raise ParseError(
                "openpyxl is required to parse XLSX files",
                file_path=str(file_path),
                suggestion="Install it with: pip install openpyxl",
            ) from None

        try:
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
        except Exception as e:
            raise ParseError(
                f"Failed to open XLSX file: {e}",
                file_path=str(file_path),
            ) from e

        sheets_text: list[str] = []
        tables: list[Table] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_data: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                row_strs = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_strs):
                    rows_data.append(row_strs)

            if not rows_data:
                continue

            # First row as headers
            headers = rows_data[0]
            data_rows = rows_data[1:]

            parts: list[str] = [f"Sheet: {sheet_name}"]
            parts.append(f"Columns: {', '.join(headers)}")
            for i, row in enumerate(data_rows):
                row_parts = []
                for j, cell in enumerate(row):
                    header = headers[j] if j < len(headers) else f"Column {j + 1}"
                    row_parts.append(f"{header}: {cell}")
                parts.append(f"Row {i + 1}: {' | '.join(row_parts)}")

            sheets_text.append("\n".join(parts))
            tables.append(Table(headers=headers, rows=data_rows, caption=sheet_name))

        wb.close()
        content = "\n\n".join(sheets_text)

        return Document(
            content=content,
            metadata=DocumentMetadata(
                source_type=SourceType.FILE,
                title=file_path.stem,
                custom={
                    "format": "xlsx",
                    "sheet_count": str(len(wb.sheetnames)),
                },
            ),
            tables=tables,
            raw_bytes=file_path.read_bytes(),
        )
